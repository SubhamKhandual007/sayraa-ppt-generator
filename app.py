print(">>> DEBUG: App file is being loaded...")
import os
import io
from flask import Flask, request, send_file, render_template, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import google.generativeai as genai
from groq import Groq
from dotenv import load_dotenv
import random

app = Flask(__name__)


# Load environment variables
load_dotenv()

# Configure APIs safely
GEMINI_KEY = os.getenv("GEMINI_API_KEY")
GROQ_KEY = os.getenv("GROQ_API_KEY")

if GEMINI_KEY:
    genai.configure(api_key=GEMINI_KEY)
    gemini_model = genai.GenerativeModel('gemini-2.0-flash')
else:
    gemini_model = None
    print("Warning: GEMINI_API_KEY not found.")

if GROQ_KEY:
    groq_client = Groq(api_key=GROQ_KEY)
else:
    groq_client = None
    print("Warning: GROQ_API_KEY not found.")

GROQ_MODEL = "llama-3.3-70b-versatile"


# Slide creation functions (unchanged from your code)
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(*tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4)))

def get_ai_theme(topic):
    try:
        prompt = f"Based on the topic '{topic}', suggest a professional color palette. Return ONLY two hex codes (primary and secondary) separated by a comma. Example: #0D47A1,#FF9800"
        
        # Try Groq first
        try:
            if not groq_client:
                raise Exception("Groq client not initialized")
            completion = groq_client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
            )
            colors = completion.choices[0].message.content.strip().split(',')
        except Exception:
            if not gemini_model:
                return ["#0D47A1", "#FF9800"]
            response = gemini_model.generate_content(prompt)
            colors = response.text.strip().split(',')
        
        if len(colors) >= 2:
            return [c.strip() for c in colors[:2]]
    except Exception:
        pass
    return ["#0D47A1", "#FF9800"] # Default colors

def generate_slide_content(section, topic, language):
    try:
        prompt = f"""
        Create minimalist, professional content for a presentation slide about '{section}' in the context of '{topic}'.
        Rules:
        - Output must be in {language}.
        - Provide exactly 2 to 3 very concise, high-impact bullet points.
        - Focus on absolute accuracy and maximum clarity.
        - Use very few words; avoid any overflow.
        - Start each line with a bullet symbol (•).
        - End with a one-sentence "Key Insight".
        """
        
        # Try Groq first
        try:
            if not groq_client:
                raise Exception("Groq client not initialized")
            completion = groq_client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
            )
            return completion.choices[0].message.content
        except Exception as groq_err:
            print(f"Groq error: {groq_err}, falling back to Gemini")
            # Fallback to Gemini
            if not gemini_model:
                return "• API keys missing.\n• Could not generate content."
            response = gemini_model.generate_content(prompt)
            return response.text
    except Exception as e:
        print(f"Error generating content for {section}: {e}")
        return f"• Content for {section} could not be generated.\n• Please add details manually."

def get_template_styles(template_name, theme_colors):
    """Define styles for different templates."""
    primary = theme_colors[0]
    secondary = theme_colors[1]
    
    styles = {
        'modern': {
            'bg': RGBColor(255, 255, 255),
            'title_color': hex_to_rgb(primary),
            'text_color': RGBColor(44, 62, 80),
            'accent': hex_to_rgb(secondary),
            'shape': MSO_SHAPE.RECTANGLE
        },
        'minimal': {
            'bg': RGBColor(255, 255, 255),
            'title_color': RGBColor(0, 0, 0),
            'text_color': RGBColor(50, 50, 50),
            'accent': RGBColor(200, 200, 200),
            'shape': None
        },
        'creative': {
            'bg': RGBColor(245, 245, 250),
            'title_color': hex_to_rgb(primary),
            'text_color': RGBColor(33, 33, 33),
            'accent': hex_to_rgb(secondary),
            'shape': MSO_SHAPE.ROUNDED_RECTANGLE
        },
        'dark': {
            'bg': RGBColor(20, 20, 25),
            'title_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(200, 200, 200),
            'accent': hex_to_rgb(primary),
            'shape': MSO_SHAPE.RECTANGLE
        }
    }
    return styles.get(template_name, styles['modern'])

def generate_sections(topic, language, num_slides):
    try:
        num_slides = int(num_slides)
        prompt = f"Generate exactly {num_slides} key sections/subtopics for a professional presentation about '{topic}'. Return as a comma-separated list in {language}."
        
        # Try Groq first
        try:
            if not groq_client:
                raise Exception("Groq client not initialized")
            completion = groq_client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
            )
            sections_text = completion.choices[0].message.content
        except Exception as groq_err:
            print(f"Groq error: {groq_err}, falling back to Gemini")
            # Fallback to Gemini
            if not gemini_model:
                return ["Introduction", "Overview", "Details"][:num_slides]
            response = gemini_model.generate_content(prompt)
            sections_text = response.text
            
        sections = [s.strip() for s in sections_text.split(',')]
        # Ensure exact count
        if len(sections) > num_slides:
            sections = sections[:num_slides]
        return sections
    except Exception as e:
        print(f"Error generating sections: {e}")
        return ["Introduction", "Overview", "Details"][:int(num_slides)]

def add_welcome_slide(prs, title_text, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, hex_to_rgb(theme[0]))
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    add_decorative_elements(slide, theme)

def add_content_page(prs, sections, topic, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    set_slide_background(slide, RGBColor(240, 240, 240))
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"Presentation Outline: {topic}"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = hex_to_rgb(theme[0])
    p.alignment = PP_ALIGN.LEFT
    for i, section in enumerate(sections):
        row = i // 3
        col = i % 3
        left_pos = Inches(0.5 + col * 3)
        top_pos = Inches(1.8 + row * 1.8)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos, Inches(2.8), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = hex_to_rgb(theme[0])
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        p = tf.add_paragraph()
        p.text = f"{i+1}. {section}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    add_watermark(slide)

def add_content_slide(prs, title_text, content_text, topic, template_style, section):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
        
    set_slide_background(slide, template_style['bg'])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = template_style['title_color']
    p.alignment = PP_ALIGN.LEFT

    # Accent Line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, Inches(1.1), width, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = template_style['accent']
    line.line.visible = False

    # Content Box (Full Width Minimalist)
    content_width = Inches(8.5)
    content_left = Inches(0.75)
    
    if template_style['shape']:
        content_box = slide.shapes.add_shape(
            template_style['shape'],
            content_left, Inches(1.8), content_width, Inches(4))
        content_box.fill.solid()
        if template_style['bg'] == RGBColor(20, 20, 25): # Dark mode
            content_box.fill.fore_color.rgb = RGBColor(40, 40, 50)
        else:
            content_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
        content_box.line.color.rgb = template_style['accent']
        content_box.line.width = Pt(1)
        tf = content_box.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        txBox = slide.shapes.add_textbox(content_left, Inches(1.8), content_width, Inches(4))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = content_text
    p.font.size = Pt(20)
    p.font.color.rgb = template_style['text_color']
    
    add_watermark(slide)

def add_thank_you_slide(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)
    set_slide_background(slide, hex_to_rgb(theme[0]))
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    add_decorative_elements(slide, theme)

def add_decorative_elements(slide, theme):
    for i in range(3):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5 + i*3), Inches(6.5), Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(theme[1])
        shape.line.color.rgb = hex_to_rgb(theme[1])

def add_watermark(slide):
    left = Inches(7.5)
    top = Inches(6.8)
    width = Inches(2.5)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Generated by Subham Khandual | SuuSri AI"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.RIGHT

def create_presentation(title, topic, language, num_slides, template_name):
    prs = Presentation()
    # Set document properties
    core_props = prs.core_properties
    core_props.author = "Subham Khandual"
    core_props.title = title
    
    theme = get_ai_theme(topic)
    style = get_template_styles(template_name, theme)
    sections = generate_sections(topic, language, num_slides)
    
    add_welcome_slide(prs, title, theme)
    add_content_page(prs, sections, topic, theme)
    
    for section in sections:
        content = generate_slide_content(section, topic, language)
        add_content_slide(prs, section, content, topic, style, section)
        
    add_thank_you_slide(prs, theme)
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Flask routes
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/favicon.ico')
def favicon():
    return send_file(os.path.join(app.root_path, 'static', 'logo-192.png'), mimetype='image/png')

@app.route('/manifest.json')
def manifest():
    return send_file(os.path.join(app.root_path, 'static', 'manifest.json'), mimetype='application/manifest+json')

@app.route('/sitemap.xml')
def sitemap():
    import datetime
    today = datetime.date.today().isoformat()
    sitemap_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">
  <url>
    <loc>{request.host_url}</loc>
    <lastmod>{today}</lastmod>
    <changefreq>weekly</changefreq>
    <priority>1.0</priority>
  </url>
</urlset>"""
    return app.response_class(sitemap_xml, mimetype='application/xml')

@app.route('/robots.txt')
def robots():
    robots_txt = f"User-agent: *\nAllow: /\nSitemap: {request.host_url}sitemap.xml\n"
    return app.response_class(robots_txt, mimetype='text/plain')

def cleanup_old_files():
    try:
        import tempfile
        import time
        temp_dir = os.path.join(tempfile.gettempdir(), 'suusri_ppt')
        if os.path.exists(temp_dir):
            now = time.time()
            for f in os.listdir(temp_dir):
                fp = os.path.join(temp_dir, f)
                if os.path.isfile(fp):
                    # Delete files older than 10 minutes (600 seconds)
                    if now - os.path.getmtime(fp) > 600:
                        try:
                            os.remove(fp)
                        except Exception:
                            pass
    except Exception as e:
        print(f"Error cleaning up old files: {e}")

@app.route('/generate', methods=['POST'])
def generate_ppt():
    if request.is_json:
        data = request.get_json()
    else:
        data = request.form
        
    title = data.get('title')
    topic = data.get('topic')
    language = data.get('language', 'English')
    template = data.get('template', 'modern')
    num_slides = 8 # Force 8 slides as requested

    # Helper function to create error response with cookie
    def make_error_response(err_msg, code=500):
        import urllib.parse
        encoded_msg = urllib.parse.quote(err_msg)
        response = jsonify({'error': err_msg})
        response.status_code = code
        response.set_cookie('ppt_error', encoded_msg, path='/')
        return response

    if not title or not topic:
        return make_error_response('Title and topic are required', 400)

    try:
        buffer = create_presentation(title, topic, language, num_slides, template)
        response = send_file(
            buffer,
            as_attachment=True,
            download_name=f"{title.replace(' ', '_')}_Presentation.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response.set_cookie('ppt_downloaded', 'true', path='/')
        return response
    except Exception as e:
        print(f"Generation error: {e}")
        return make_error_response(str(e), 500)

if __name__ == '__main__':
    app.run(debug=True, port=5000)